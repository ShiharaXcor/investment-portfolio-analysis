from pathlib import Path
from typing import Tuple, List, Dict, Any
import io, re
import pandas as pd
from pypdf import PdfReader
import pdfplumber
from docx import Document as Docx
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract

def _page(page:int, text:str, section:str|None=None) -> Dict[str, Any]:
    return {"page": page, "text": text or "", "section": section}

def _pdf_native(fp:str) -> List[Dict[str, Any]]:
    reader = PdfReader(fp)
    pages = []
    has_text = False
    for i, _ in enumerate(reader.pages):
        t = reader.pages[i].extract_text() or ""
        has_text = has_text or bool(t.strip())
        pages.append(_page(i+1, t))
    return pages, has_text

def _pdf_ocr(fp:str) -> List[Dict[str, Any]]:
    pages = []
    with pdfplumber.open(fp) as pdf:
        for i, p in enumerate(pdf.pages):
            img = p.to_image(resolution=300).original
            buf = io.BytesIO()
            img.save(buf, format="PNG"); buf.seek(0)
            text = pytesseract.image_to_string(Image.open(buf)) or ""
            pages.append(_page(i+1, text))
    return pages

def extract_pdf(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    pages, has_text = _pdf_native(fp)
    if not has_text:
        pages = _pdf_ocr(fp)
    full = "\n".join([x["text"] for x in pages])
    return full, pages

def extract_docx(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    doc = Docx(fp)
    text = "\n".join(p.text for p in doc.paragraphs)
    return text, [_page(1, text)]

def extract_tabular(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    p = Path(fp)
    if p.suffix.lower() == ".csv":
        df = pd.read_csv(fp)
    else:
        df = pd.read_excel(fp)
    text = df.to_csv(index=False)
    return text, [_page(1, text, "Table")]

def extract_pptx(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    prs = Presentation(fp)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        out.append("\n".join(slide_text))
    full = "\n\n".join(out)
    pages = [_page(i+1, t, f"Slide {i+1}") for i, t in enumerate(out)]
    return full, pages

def extract_html(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    html = Path(fp).read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for s in soup(["script","style","noscript"]): s.decompose()
    text = soup.get_text("\n")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text, [_page(1, text, "HTML")]

def extract_txt(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    text = Path(fp).read_text(encoding="utf-8", errors="ignore")
    return text, [_page(1, text)]

def extract_image(fp:str) -> Tuple[str, List[Dict[str, Any]]]:
    img = Image.open(fp)
    text = pytesseract.image_to_string(img) or ""
    return text, [_page(1, text, "Image OCR")]

def extract_any(fp:str) -> Tuple[str, Dict[str, Any], List[Dict[str, Any]]]:
    p = Path(fp)
    name = p.name
    ext = p.suffix.lower()
    if ext == ".pdf":
        full, pages = extract_pdf(fp)
    elif ext == ".docx":
        full, pages = extract_docx(fp)
    elif ext in {".csv", ".xlsx"}:
        full, pages = extract_tabular(fp)
    elif ext == ".pptx":
        full, pages = extract_pptx(fp)
    elif ext in {".html", ".htm"}:
        full, pages = extract_html(fp)
    elif ext in {".txt", ".md"}:
        full, pages = extract_txt(fp)
    elif ext in {".png",".jpg",".jpeg",".webp",".tiff"}:
        full, pages = extract_image(fp)
    else:
        full, pages = extract_txt(fp)
    meta = {"document_name": name}
    return full, meta, pages
